import os
import chardet
import pandas as pd
from pypdf import PdfReader
from docx import Document
from bs4 import BeautifulSoup

# ── Optional imports (gracefully degrade if not installed) ──────────────────
try:
    from pptx import Presentation as PptxPresentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    _RTF_AVAILABLE = True
except ImportError:
    _RTF_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image as PilImage
    import sys
    if sys.platform == "win32":
        pytesseract.pytesseract.tesseract_cmd = r"C:\Users\KIIT0001\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"
    _OCR_AVAILABLE = True
except ImportError:
    _OCR_AVAILABLE = False

try:
    from PIL import Image as PilImage
    from PIL.ExifTags import TAGS
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

# ── File-type groups ────────────────────────────────────────────────────────
IMAGE_EXTS      = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".gif", ".webp"}
PPTX_EXTS       = {".pptx", ".ppt"}
SPREADSHEET_EXTS= {".xlsx", ".xls", ".csv"}
WORD_EXTS       = {".docx", ".doc"}
HTML_EXTS       = {".html", ".htm"}
XML_EXTS        = {".xml"}
RTF_EXTS        = {".rtf"}
PLAIN_TEXT_EXTS = {
    ".txt", ".log", ".json", ".md", ".markdown",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".py",  ".js",  ".ts",  ".jsx", ".tsx", ".java",
    ".c",   ".cpp", ".h",   ".cs",  ".go",  ".rs",
    ".sh",  ".bat", ".ps1", ".sql", ".r",   ".rb",
    ".env", ".gitignore", ".dockerfile",
}


def _read_text_file(file_path: str) -> str:
    """Read any text-based file with automatic charset detection."""
    with open(file_path, "rb") as f:
        raw = f.read()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"
    return raw.decode(encoding, errors="replace")


def _extract_image(file_path: str) -> str:
    """
    Attempt OCR with pytesseract; fall back to Pillow image metadata
    (dimensions, mode, format, EXIF) so images are never silently bypassed.
    """
    parts = []

    if _OCR_AVAILABLE:
        try:
            img = PilImage.open(file_path)
            
            # Preprocessing for better OCR on UI screenshots / dark mode
            # 1. Upscale by 2x for clearer small fonts
            img = img.resize((img.width * 2, img.height * 2), PilImage.Resampling.LANCZOS)
            # 2. Convert to grayscale
            img = img.convert('L')
            
            # 3. Let Tesseract use default Page Segmentation Mode (PSM 3) which handles complex documents better
            ocr_text = pytesseract.image_to_string(img).strip()
            
            if ocr_text:
                parts.append(f"[OCR Text]\n{ocr_text}")
        except Exception as e:
            parts.append(f"[OCR Error]: {e}")

    if _PIL_AVAILABLE:
        try:
            img = PilImage.open(file_path)
            parts.append(
                f"[Image Metadata]\n"
                f"  Format : {img.format}\n"
                f"  Mode   : {img.mode}\n"
                f"  Size   : {img.size[0]} x {img.size[1]} pixels"
            )
            # Extract human-readable EXIF tags if present
            exif_data = img._getexif() if hasattr(img, "_getexif") else None
            if exif_data:
                exif_lines = []
                for tag_id, value in exif_data.items():
                    tag_name = TAGS.get(tag_id, tag_id)
                    # Skip raw binary blobs
                    if isinstance(value, bytes):
                        continue
                    exif_lines.append(f"  {tag_name}: {value}")
                if exif_lines:
                    parts.append("[EXIF Data]\n" + "\n".join(exif_lines))
        except Exception as e:
            parts.append(f"[Image Read Error]: {e}")

    if not parts:
        return f"[Image]: Could not extract content. Install pytesseract + Tesseract for OCR."

    return "\n\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    """Extract all text from PowerPoint slides."""
    prs = PptxPresentation(file_path)
    slide_texts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slide_texts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slide_texts)


def _extract_html(file_path: str) -> str:
    """Strip HTML tags and return clean text."""
    raw = _read_text_file(file_path)
    soup = BeautifulSoup(raw, "lxml")
    return soup.get_text(separator="\n", strip=True)


def _extract_xml(file_path: str) -> str:
    """Return pretty-printed XML content (BeautifulSoup preserves structure)."""
    raw = _read_text_file(file_path)
    soup = BeautifulSoup(raw, "lxml-xml")
    return soup.get_text(separator="\n", strip=True)


def _try_as_text(file_path: str) -> str:
    """
    Final fallback: attempt to read any unknown file as text with
    charset auto-detection. If it's truly binary, return a helpful message.
    """
    try:
        text = _read_text_file(file_path)
        # Heuristic: if >30 % of chars are non-printable it's binary
        non_printable = sum(1 for c in text if not c.isprintable() and c not in "\n\r\t")
        if len(text) > 0 and non_printable / len(text) > 0.30:
            return f"[Binary File]: Content is not human-readable text (detected as binary)."
        return text
    except Exception as e:
        return f"[Read Error]: {e}"


# ── Public API ───────────────────────────────────────────────────────────────

def extract_text_from_file(file_path: str) -> str:
    """
    Inspects a file's extension and routes it to the correct extraction
    handler, returning a unified text string. Supports:
      PDF, DOCX/DOC, XLSX/XLS/CSV, PPTX/PPT,
      HTML/HTM, XML, RTF,
      PNG/JPG/JPEG/BMP/TIFF/GIF/WEBP (OCR + metadata),
      TXT/LOG/JSON/MD/YAML/TOML/INI/CFG and most source-code formats,
      plus a smart binary-safe fallback for everything else.
    """
    ext = os.path.splitext(file_path)[1].lower()
    extracted_text = ""

    try:
        # ── 1. PDF ────────────────────────────────────────────────────────
        if ext == ".pdf":
            try:
                import fitz
                doc = fitz.open(file_path)
                slices = []
                for page in doc:
                    page_text = page.get_text()
                    if len(page_text.strip()) < 50 and _OCR_AVAILABLE:
                        # Likely a scanned page, fallback to OCR
                        pix = page.get_pixmap(dpi=150)
                        img = PilImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        page_text = pytesseract.image_to_string(img)
                    slices.append(page_text)
                extracted_text = "\n".join(slices)
            except ImportError:
                # Fallback if PyMuPDF isn't installed
                reader = PdfReader(file_path)
                slices = [p.extract_text(extraction_mode="layout") for p in reader.pages]
                slices = [s for s in slices if s]
                extracted_text = "\n".join(slices)

        # ── 2. Word ───────────────────────────────────────────────────────
        elif ext in WORD_EXTS:
            doc = Document(file_path)
            slices = [p.text for p in doc.paragraphs if p.text]
            extracted_text = "\n".join(slices)

        # ── 3. Spreadsheets ───────────────────────────────────────────────
        elif ext in SPREADSHEET_EXTS:
            df = pd.read_csv(file_path) if ext == ".csv" else pd.read_excel(file_path)
            extracted_text = df.to_string(index=False)

        # ── 4. PowerPoint ─────────────────────────────────────────────────
        elif ext in PPTX_EXTS:
            if _PPTX_AVAILABLE:
                extracted_text = _extract_pptx(file_path)
            else:
                extracted_text = "[PPTX Error]: python-pptx not installed. Run: pip install python-pptx"

        # ── 5. HTML ───────────────────────────────────────────────────────
        elif ext in HTML_EXTS:
            extracted_text = _extract_html(file_path)

        # ── 6. XML ────────────────────────────────────────────────────────
        elif ext in XML_EXTS:
            extracted_text = _extract_xml(file_path)

        # ── 7. RTF ────────────────────────────────────────────────────────
        elif ext in RTF_EXTS:
            if _RTF_AVAILABLE:
                raw = _read_text_file(file_path)
                extracted_text = rtf_to_text(raw)
            else:
                extracted_text = "[RTF Error]: striprtf not installed. Run: pip install striprtf"

        # ── 8. Images (OCR + metadata) ────────────────────────────────────
        elif ext in IMAGE_EXTS:
            extracted_text = _extract_image(file_path)

        # ── 9. Known plain-text / code formats ───────────────────────────
        elif ext in PLAIN_TEXT_EXTS:
            extracted_text = _read_text_file(file_path)

        # ── 10. Unknown — smart binary-safe fallback ──────────────────────
        else:
            extracted_text = _try_as_text(file_path)

        # Sanitize for PostgreSQL: Remove null bytes (\x00)
        if isinstance(extracted_text, str):
            extracted_text = extracted_text.replace('\x00', '')

        return extracted_text.strip()

    except Exception as e:
        return f"[Parsing Failure]: Could not extract text. Details: {str(e)}"